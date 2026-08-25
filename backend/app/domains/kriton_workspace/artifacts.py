from __future__ import annotations

import asyncio
import hashlib
import io
import re
import uuid
from datetime import datetime, timezone
from pathlib import Path

from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.domains.kriton_workspace.models import GeneratedArtifact
from app.core.config import get_settings
from app.domains.kriton_workspace.storage import download_object, upload_object, uses_supabase
from app.domains.kriton_workspace.retention import artifact_expiry
from app.domains.kriton_workspace.document_spec import (
    DocumentBlock, DocumentSpec, build_document_spec, normalize_document_text,
    visual_as_data_uri,
)

_ROOT = Path(__file__).resolve().parents[3] / "data" / "workspace_artifacts"
_SAFE = re.compile(r"[^A-Za-z0-9._-]+")


def _plain(text: str) -> str:
    text = re.sub(r"<br\s*/?>", "\n", text, flags=re.IGNORECASE)
    text = re.sub(r"<[^>]+>", "", text)
    text = re.sub(r"[*_`]", "", re.sub(r"\[([^]]+)]\([^)]+\)", r"\1", text))
    return normalize_document_text(text)


def _docx_bytes(spec: DocumentSpec) -> bytes:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches
    document = Document()
    document.add_heading(spec.title, 0)
    for block in spec.blocks:
        if block.type == "heading":
            document.add_heading(_plain(block.text), level=min(block.level, 2))
        elif block.type == "paragraph":
            document.add_paragraph(_plain(block.text))
        elif block.type == "bullets":
            for item in block.items:
                document.add_paragraph(_plain(item), style="List Bullet")
        elif block.type == "table":
            table = document.add_table(rows=1, cols=len(block.headers))
            table.style = "Table Grid"
            for index, header in enumerate(block.headers):
                table.rows[0].cells[index].text = _plain(header)
            for values in block.rows:
                cells = table.add_row().cells
                for index, value in enumerate(values):
                    cells[index].text = _plain(value)
                    if re.search(r"\d", value):
                        cells[index].paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.RIGHT
        elif block.type in {"chart", "flow_diagram"} and block.image_png:
            document.add_heading(block.title, level=2)
            paragraph = document.add_paragraph()
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
            paragraph.add_run().add_picture(io.BytesIO(block.image_png), width=Inches(6.4))
    if spec.source_locations:
        document.add_heading("Source References", level=1)
        for location in spec.source_locations:
            document.add_paragraph(location, style="List Bullet")
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _xlsx_bytes(spec: DocumentSpec, analysis: dict) -> bytes:
    from openpyxl import Workbook
    from openpyxl.chart import BarChart, Reference
    from openpyxl.drawing.image import Image as WorksheetImage
    from openpyxl.styles import Font
    workbook = Workbook()
    report = workbook.active
    report.title = "Report"
    report.append([spec.title])
    report["A1"].font = Font(bold=True, size=16)
    for block in spec.blocks:
        if block.type in {"heading", "paragraph"}:
            report.append([_plain(block.text)])
            if block.type == "heading":
                report.cell(report.max_row, 1).font = Font(bold=True, size=13)
        elif block.type == "bullets":
            for item in block.items:
                report.append([f"• {_plain(item)}"])
        elif block.type == "table":
            report.append(block.headers)
            for cell in report[report.max_row]:
                cell.font = Font(bold=True)
            for row in block.rows:
                report.append(row)
        elif block.type in {"chart", "flow_diagram"} and block.image_png:
            image = WorksheetImage(io.BytesIO(block.image_png))
            image.width, image.height = 720, 405 if block.type == "chart" else 216
            anchor = f"A{report.max_row + 2}"
            report.add_image(image, anchor)
            report.append([block.title])
            for _ in range(22 if block.type == "chart" else 13):
                report.append([])
    report.column_dimensions["A"].width = 110
    kpis = workbook.create_sheet("Verified KPIs")
    kpis.append(["Metric", "Value"])
    for key, value in analysis.get("metrics", {}).items():
        kpis.append([key.replace("_", " ").title(), value])
    customers = workbook.create_sheet("Customer Analysis")
    customers.append(["Customer", "Revenue"])
    for key, value in analysis.get("customer_revenue", {}).items():
        customers.append([key, value])
    products = workbook.create_sheet("Product Analysis")
    products.append(["Product", "Revenue"])
    for key, value in analysis.get("product_revenue", {}).items():
        products.append([key, value])
    if products.max_row > 1:
        chart = BarChart()
        chart.title = "Revenue by Product"
        chart.y_axis.title = "Revenue"
        chart.add_data(Reference(products, min_col=2, min_row=1, max_row=products.max_row), titles_from_data=True)
        chart.set_categories(Reference(products, min_col=1, min_row=2, max_row=products.max_row))
        products.add_chart(chart, "D2")
    sources = workbook.create_sheet("Source References")
    sources.append(["Source location"])
    for location in analysis.get("evidence_locations", []):
        sources.append([location])
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _pdf_bytes(spec: DocumentSpec) -> bytes:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.utils import ImageReader
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image as RLImage, PageBreak
    from xml.sax.saxutils import escape
    buffer = io.BytesIO()
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name="KritonBody", parent=styles["BodyText"], leading=15, spaceAfter=8))
    styles.add(ParagraphStyle(
        name="KritonTableHeader", parent=styles["BodyText"],
        textColor=colors.white, fontName="Helvetica-Bold", leading=12,
    ))
    story = [Paragraph(escape(_plain(spec.title)), styles["Title"]), Spacer(1, 12)]
    for block in spec.blocks:
        if block.type == "heading":
            story.append(Paragraph(escape(_plain(block.text)), styles["Heading1" if block.level <= 2 else "Heading2"]))
        elif block.type == "paragraph":
            story.append(Paragraph(escape(_plain(block.text)).replace("\n", "<br/>"), styles["KritonBody"]))
        elif block.type == "bullets":
            for item in block.items:
                story.append(Paragraph(f"• {escape(_plain(item))}", styles["KritonBody"]))
        elif block.type == "table":
            data = [[Paragraph(escape(_plain(value)).replace("\n", "<br/>"), styles["KritonTableHeader"]) for value in block.headers]]
            data.extend([[Paragraph(escape(_plain(value)).replace("\n", "<br/>"), styles["BodyText"]) for value in row] for row in block.rows])
            available = A4[0] - 72
            table = Table(data, repeatRows=1, colWidths=[available / len(block.headers)] * len(block.headers))
            table.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#16799A")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#B0BEC5")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F4F8F9")]),
                ("LEFTPADDING", (0, 0), (-1, -1), 5),
                ("RIGHTPADDING", (0, 0), (-1, -1), 5),
            ] + [
                ("ALIGN", (column, 1), (column, -1), "RIGHT")
                for column in range(len(block.headers))
                if any(re.search(r"\d", row[column]) for row in block.rows)
            ]))
            story.extend([table, Spacer(1, 12)])
        elif block.type in {"chart", "flow_diagram"} and block.image_png:
            story.append(Paragraph(escape(block.title), styles["Heading2"]))
            image = RLImage(io.BytesIO(block.image_png))
            image._restrictSize(A4[0] - 90, 360)
            story.extend([image, Spacer(1, 12)])
    if spec.source_locations:
        story.append(Paragraph("Source References", styles["Heading1"]))
        for location in spec.source_locations:
            story.append(Paragraph(f"• {escape(_plain(location))}", styles["KritonBody"]))
    document = SimpleDocTemplate(buffer, pagesize=A4, rightMargin=36, leftMargin=36, topMargin=42, bottomMargin=42)
    document.build(story)
    return buffer.getvalue()


def _pptx_bytes(spec: DocumentSpec) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = spec.title
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = "Generated by Kriton from verified document evidence"

    current_title = "Report Summary"
    for block in spec.blocks:
        if block.type == "heading":
            current_title = _plain(block.text)
            continue
        if block.type in {"chart", "flow_diagram"} and block.image_png:
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = block.title
            slide.shapes.add_picture(io.BytesIO(block.image_png), Inches(0.7), Inches(1.4), width=Inches(8.6))
            continue
        if block.type == "table":
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])
            slide.shapes.title.text = current_title
            visible_rows = block.rows[:12]
            shape = slide.shapes.add_table(len(visible_rows) + 1, len(block.headers), Inches(0.45), Inches(1.35), Inches(9.1), Inches(5.5))
            table = shape.table
            for column, header in enumerate(block.headers):
                table.cell(0, column).text = _plain(header)
            for row_index, values in enumerate(visible_rows, start=1):
                for column, value in enumerate(values):
                    table.cell(row_index, column).text = _plain(value)
            continue
        lines = block.items if block.type == "bullets" else [_plain(block.text)]
        for offset in range(0, len(lines), 8):
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = current_title if offset == 0 else f"{current_title} (continued)"
            frame = slide.placeholders[1].text_frame
            frame.clear()
            for index, line in enumerate(lines[offset:offset + 8]):
                paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                paragraph.text = line[:500]
                paragraph.level = 0
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _markdown_bytes(spec: DocumentSpec) -> bytes:
    lines = [f"# {spec.title}", ""]
    for block in spec.blocks:
        if block.type == "heading":
            lines.extend([f"{'#' * min(block.level + 1, 4)} {block.text}", ""])
        elif block.type == "paragraph":
            lines.extend([block.text, ""])
        elif block.type == "bullets":
            lines.extend([*(f"- {item}" for item in block.items), ""])
        elif block.type == "table":
            lines.extend([
                "| " + " | ".join(block.headers) + " |",
                "| " + " | ".join("---" for _ in block.headers) + " |",
                *("| " + " | ".join(row) + " |" for row in block.rows),
                "",
            ])
        elif block.type in {"chart", "flow_diagram"}:
            lines.extend([f"## {block.title}", "", f"![{block.alt_text}]({visual_as_data_uri(block)})", ""])
    if spec.source_locations:
        lines.extend(["## Source References", "", *(f"- {location}" for location in spec.source_locations), ""])
    return "\n".join(lines).encode("utf-8")


def _render(format_name: str, title: str, narrative: str, analysis: dict, request_text: str = "") -> tuple[bytes, str, str]:
    spec = build_document_spec(title, narrative, analysis, request_text)
    if format_name == "xlsx":
        content = _xlsx_bytes(spec, analysis)
        from openpyxl import load_workbook
        workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=False)
        try:
            required = {"Report", "Verified KPIs", "Customer Analysis", "Product Analysis", "Source References"}
            if not required.issubset(workbook.sheetnames):
                raise ValueError("Generated XLSX failed sheet validation")
        finally:
            workbook.close()
        return content, "xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    if format_name == "pdf":
        content = _pdf_bytes(spec)
        from pypdf import PdfReader
        if len(PdfReader(io.BytesIO(content)).pages) < 1:
            raise ValueError("Generated PDF has no pages")
        return content, "pdf", "application/pdf"
    if format_name == "pptx":
        content = _pptx_bytes(spec)
        from pptx import Presentation
        if len(Presentation(io.BytesIO(content)).slides) < 2:
            raise ValueError("Generated PowerPoint has no content slides")
        return content, "pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    if format_name == "md":
        content = _markdown_bytes(spec)
        if len(content.strip()) < 10:
            raise ValueError("Generated Markdown is empty")
        return content, "md", "text/markdown; charset=utf-8"
    content = _docx_bytes(spec)
    from docx import Document
    if not Document(io.BytesIO(content)).paragraphs:
        raise ValueError("Generated DOCX has no readable paragraphs")
    return content, "docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"


async def create_generated_artifact(
    db: AsyncSession, *, title: str, narrative: str, analysis: dict, format_name: str,
    tenant_id: str, user_id: str, conversation_id: str | None, query_id: str,
    source_document_ids: list[str], request_text: str = "",
) -> GeneratedArtifact:
    content, extension, mime_type = await asyncio.to_thread(
        _render, format_name, title, narrative, analysis, request_text
    )
    artifact_id = f"artifact-{uuid.uuid4().hex[:16]}"
    filename = f"{_SAFE.sub('_', title).strip('_') or 'Kriton_Report'}.{extension}"
    object_path = f"{tenant_id}/{user_id}/{conversation_id or 'no-conversation'}/{artifact_id}/{filename}"
    provider = "supabase" if uses_supabase() else "local"
    bucket = get_settings().ARTIFACT_STORAGE_BUCKET if provider == "supabase" else None
    if provider == "supabase":
        await upload_object(bucket or "", object_path, content, mime_type)
        stored_path = object_path
    else:
        tenant_dir = _ROOT / _SAFE.sub("_", tenant_id)
        tenant_dir.mkdir(parents=True, exist_ok=True)
        path = tenant_dir / f"{artifact_id}_{filename}"
        path.write_bytes(content)
        stored_path = str(path.relative_to(_ROOT.parents[1]))
    artifact = GeneratedArtifact(
        id=artifact_id, tenant_id=tenant_id, user_id=user_id,
        conversation_id=conversation_id, query_id=query_id,
        source_document_ids=source_document_ids, filename=filename,
        mime_type=mime_type,
        storage_path=stored_path,
        storage_provider=provider,
        storage_bucket=bucket,
        storage_object_path=object_path if provider == "supabase" else None,
        content_hash=hashlib.sha256(content).hexdigest(), status="READY",
        expires_at=artifact_expiry(),
    )
    db.add(artifact)
    await db.commit()
    await db.refresh(artifact)
    return artifact


async def get_generated_artifact(db: AsyncSession, *, artifact_id: str, tenant_id: str, user_id: str) -> GeneratedArtifact | None:
    result = await db.execute(select(GeneratedArtifact).where(
        GeneratedArtifact.id == artifact_id,
        GeneratedArtifact.tenant_id == tenant_id,
        GeneratedArtifact.user_id == user_id,
        GeneratedArtifact.status == "READY",
        GeneratedArtifact.expires_at > datetime.now(timezone.utc),
    ))
    return result.scalar_one_or_none()


def artifact_absolute_path(artifact: GeneratedArtifact) -> Path:
    return (_ROOT.parents[1] / artifact.storage_path).resolve()


async def artifact_bytes(artifact: GeneratedArtifact) -> bytes:
    if artifact.storage_provider == "supabase" and artifact.storage_bucket and artifact.storage_object_path:
        return await download_object(artifact.storage_bucket, artifact.storage_object_path)
    return artifact_absolute_path(artifact).read_bytes()
