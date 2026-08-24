from __future__ import annotations

import asyncio
import hashlib
import io
import re
import uuid
import zipfile
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from xml.etree import ElementTree

from sqlalchemy import delete, func, select, text
from sqlalchemy.ext.asyncio import AsyncSession

from app.domains.kriton_workspace.models import ConversationDocument, Document, DocumentChunk
from app.orchestration.websearch import WebSource
from app.core.config import get_settings
from app.domains.kriton_workspace.storage import delete_object, download_object, upload_object, uses_supabase
from app.domains.kriton_workspace.retention import document_expiry, failed_upload_expiry

_STORAGE_ROOT = Path(__file__).resolve().parents[3] / "data" / "workspace_documents"
_SAFE_NAME = re.compile(r"[^A-Za-z0-9._-]+")
_WORD = re.compile(r"[A-Za-z0-9][A-Za-z0-9_.%/-]*")
_STOP_WORDS = {
    "a", "an", "and", "are", "as", "at", "be", "by", "for", "from", "how",
    "in", "is", "it", "of", "on", "or", "that", "the", "this", "to", "was",
    "what", "when", "where", "which", "who", "why", "with",
}
_MAX_CHARS = 3500
_OVERLAP = 250


@dataclass(frozen=True)
class ExtractedBlock:
    text: str
    location: str
    metadata: dict


def _chunk_block(block: ExtractedBlock) -> list[ExtractedBlock]:
    # Preserve tabs emitted by table/spreadsheet processors so column
    # boundaries remain available to retrieval and the answering model.
    text = re.sub(r" +", " ", block.text).strip()
    if not text:
        return []
    if len(text) <= _MAX_CHARS:
        return [ExtractedBlock(text, block.location, block.metadata)]

    chunks: list[ExtractedBlock] = []
    start = 0
    while start < len(text):
        end = min(start + _MAX_CHARS, len(text))
        if end < len(text):
            boundary = text.rfind("\n", start, end)
            if boundary < start + (_MAX_CHARS // 2):
                boundary = text.rfind(". ", start, end)
            if boundary > start:
                end = boundary + 1
        chunks.append(ExtractedBlock(text[start:end].strip(), block.location, block.metadata))
        if end >= len(text):
            break
        start = max(end - _OVERLAP, start + 1)
    return chunks


def _extract_pdf(content: bytes) -> list[ExtractedBlock]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    return [
        ExtractedBlock(page.extract_text() or "", f"Page {index}", {"page": index})
        for index, page in enumerate(reader.pages, start=1)
    ]


def _extract_docx(content: bytes) -> list[ExtractedBlock]:
    from docx import Document as DocxDocument

    document = DocxDocument(io.BytesIO(content))
    blocks: list[ExtractedBlock] = []
    paragraphs = "\n".join(p.text for p in document.paragraphs if p.text.strip())
    if paragraphs:
        blocks.append(ExtractedBlock(paragraphs, "Document text", {"kind": "text"}))
    for index, table in enumerate(document.tables, start=1):
        rows = ["\t".join(cell.text.strip() for cell in row.cells) for row in table.rows]
        blocks.append(ExtractedBlock("\n".join(rows), f"Table {index}", {"kind": "table", "table": index}))
    return blocks


def _extract_xlsx_openpyxl(content: bytes) -> list[ExtractedBlock]:
    from openpyxl import load_workbook

    formulas = load_workbook(io.BytesIO(content), read_only=True, data_only=False)
    values = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    blocks: list[ExtractedBlock] = []
    try:
        for sheet in formulas.worksheets:
            value_sheet = values[sheet.title]
            batch: list[str] = []
            start_row = 1
            for row_number, (formula_row, value_row) in enumerate(
                zip(sheet.iter_rows(), value_sheet.iter_rows()), start=1
            ):
                rendered: list[str] = []
                for formula_cell, value_cell in zip(formula_row, value_row):
                    value = value_cell.value
                    formula = formula_cell.value
                    if formula is None and value is None:
                        rendered.append("")
                    elif isinstance(formula, str) and formula.startswith("="):
                        rendered.append(f"{formula} => {value}")
                    else:
                        rendered.append(str(value if value is not None else formula))
                if any(cell for cell in rendered):
                    if not batch:
                        start_row = row_number
                    batch.append("\t".join(rendered).rstrip("\t"))
                if len(batch) >= 40:
                    location = f"'{sheet.title}'!{start_row}:{row_number}"
                    blocks.append(ExtractedBlock("\n".join(batch), location, {"sheet": sheet.title, "start_row": start_row, "end_row": row_number}))
                    batch = []
            if batch:
                end_row = start_row + len(batch) - 1
                location = f"'{sheet.title}'!{start_row}:{end_row}"
                blocks.append(ExtractedBlock("\n".join(batch), location, {"sheet": sheet.title, "start_row": start_row, "end_row": end_row}))
    finally:
        formulas.close()
        values.close()
    return blocks


def _xlsx_column_number(reference: str) -> int:
    letters = "".join(character for character in reference if character.isalpha()).upper()
    value = 0
    for letter in letters:
        value = value * 26 + ord(letter) - 64
    return value


def _extract_xlsx_zip(content: bytes) -> list[ExtractedBlock]:
    """Dependency-free XLSX fallback used when openpyxl is unavailable."""
    ns = {"m": "http://schemas.openxmlformats.org/spreadsheetml/2006/main"}
    rel_ns = {"r": "http://schemas.openxmlformats.org/package/2006/relationships"}
    office_rel = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    with zipfile.ZipFile(io.BytesIO(content)) as archive:
        shared: list[str] = []
        if "xl/sharedStrings.xml" in archive.namelist():
            root = ElementTree.fromstring(archive.read("xl/sharedStrings.xml"))
            shared = ["".join(node.text or "" for node in item.findall(".//m:t", ns)) for item in root.findall("m:si", ns)]

        workbook = ElementTree.fromstring(archive.read("xl/workbook.xml"))
        relationships = ElementTree.fromstring(archive.read("xl/_rels/workbook.xml.rels"))
        targets = {rel.attrib["Id"]: rel.attrib["Target"] for rel in relationships.findall("r:Relationship", rel_ns)}
        blocks: list[ExtractedBlock] = []
        for sheet in workbook.findall(".//m:sheet", ns):
            title = sheet.attrib["name"]
            target = targets[sheet.attrib[office_rel]].lstrip("/")
            path = target if target.startswith("xl/") else f"xl/{target}"
            root = ElementTree.fromstring(archive.read(path))
            batch: list[str] = []
            start_row = 1
            last_row = 1
            for row in root.findall(".//m:sheetData/m:row", ns):
                row_number = int(row.attrib.get("r", last_row))
                rendered: list[str] = []
                for cell in row.findall("m:c", ns):
                    column = _xlsx_column_number(cell.attrib.get("r", "A1"))
                    while len(rendered) < column - 1:
                        rendered.append("")
                    value_node = cell.find("m:v", ns)
                    formula_node = cell.find("m:f", ns)
                    raw = value_node.text if value_node is not None else ""
                    if cell.attrib.get("t") == "s" and raw:
                        raw = shared[int(raw)]
                    elif cell.attrib.get("t") == "inlineStr":
                        raw = "".join(node.text or "" for node in cell.findall(".//m:t", ns))
                    rendered.append(f"={formula_node.text} => {raw}" if formula_node is not None else raw)
                if any(rendered):
                    if not batch:
                        start_row = row_number
                    batch.append("\t".join(rendered).rstrip("\t"))
                    last_row = row_number
                if len(batch) >= 40:
                    location = f"'{title}'!{start_row}:{last_row}"
                    blocks.append(ExtractedBlock("\n".join(batch), location, {"sheet": title, "start_row": start_row, "end_row": last_row}))
                    batch = []
            if batch:
                location = f"'{title}'!{start_row}:{last_row}"
                blocks.append(ExtractedBlock("\n".join(batch), location, {"sheet": title, "start_row": start_row, "end_row": last_row}))
        return blocks


def _extract_xlsx(content: bytes) -> list[ExtractedBlock]:
    try:
        return _extract_xlsx_openpyxl(content)
    except ModuleNotFoundError:
        return _extract_xlsx_zip(content)


def _extract_pptx(content: bytes) -> list[ExtractedBlock]:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(content))
    blocks: list[ExtractedBlock] = []
    for index, slide in enumerate(presentation.slides, start=1):
        text = "\n".join(
            shape.text for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        )
        blocks.append(ExtractedBlock(text, f"Slide {index}", {"slide": index}))
    return blocks


def extract_document(filename: str, content: bytes) -> list[ExtractedBlock]:
    suffix = Path(filename).suffix.lower()
    extractor = {
        ".pdf": _extract_pdf,
        ".docx": _extract_docx,
        ".xlsx": _extract_xlsx,
        ".pptx": _extract_pptx,
    }.get(suffix)
    if extractor is None:
        raise ValueError(f"Unsupported file type: {suffix or 'unknown'}")
    return [chunk for block in extractor(content) for chunk in _chunk_block(block)]


async def create_document_upload(
    db: AsyncSession,
    *,
    filename: str,
    mime_type: str,
    content: bytes,
    tenant_id: str,
    user_id: str,
) -> Document:
    document_id = f"doc-{uuid.uuid4().hex[:16]}"
    safe_name = _SAFE_NAME.sub("_", filename)
    object_path = f"{tenant_id}/{user_id}/{document_id}/{safe_name}"
    provider = "supabase" if uses_supabase() else "local"
    bucket = get_settings().DOCUMENT_STORAGE_BUCKET if provider == "supabase" else None
    if provider == "supabase":
        await upload_object(bucket or "", object_path, content, mime_type)
        stored_path = object_path
    else:
        tenant_dir = _STORAGE_ROOT / _SAFE_NAME.sub("_", tenant_id)
        tenant_dir.mkdir(parents=True, exist_ok=True)
        path = tenant_dir / f"{document_id}_{safe_name}"
        path.write_bytes(content)
        stored_path = str(path.relative_to(_STORAGE_ROOT.parents[1]))

    document = Document(
        id=document_id,
        tenant_id=tenant_id,
        user_id=user_id,
        filename=filename,
        mime_type=mime_type,
        storage_path=stored_path,
        storage_provider=provider,
        storage_bucket=bucket,
        storage_object_path=object_path if provider == "supabase" else None,
        content_hash=hashlib.sha256(content).hexdigest(),
        status="UPLOADED",
        expires_at=document_expiry(),
    )
    db.add(document)
    await db.commit()
    await db.refresh(document)
    return document


async def process_document(db: AsyncSession, document: Document) -> int:
    """Extract and index one stored document; safe to retry after failure."""
    document.status = "PROCESSING"
    document.processing_error = None
    await db.execute(delete(DocumentChunk).where(DocumentChunk.document_id == document.id))
    await db.commit()
    try:
        if document.storage_provider == "supabase":
            content = await download_object(
                document.storage_bucket or "", document.storage_object_path or document.storage_path
            )
        else:
            storage_path = (_STORAGE_ROOT.parents[1] / document.storage_path).resolve()
            if not storage_path.is_relative_to(_STORAGE_ROOT.resolve()):
                raise ValueError("Document storage path escapes the configured root")
            content = await asyncio.to_thread(storage_path.read_bytes)
        blocks = await asyncio.to_thread(extract_document, document.filename, content)
        if not blocks:
            raise ValueError("No readable text or cells were found in the document")
        db.add_all([
            DocumentChunk(
                document_id=document.id,
                tenant_id=document.tenant_id,
                ordinal=index,
                text=block.text,
                location=block.location,
                chunk_metadata=block.metadata,
            )
            for index, block in enumerate(blocks)
        ])
        document.status = "READY"
    except Exception as exc:
        document.status = "FAILED"
        document.processing_error = str(exc)[:1000]
        document.expires_at = failed_upload_expiry()
    await db.commit()
    await db.refresh(document)
    return len(blocks) if document.status == "READY" else 0


async def ingest_document(
    db: AsyncSession,
    *,
    filename: str,
    mime_type: str,
    content: bytes,
    tenant_id: str,
    user_id: str,
) -> tuple[Document, int]:
    """Synchronous compatibility path used by local development and tests."""
    document = await create_document_upload(
        db,
        filename=filename,
        mime_type=mime_type,
        content=content,
        tenant_id=tenant_id,
        user_id=user_id,
    )
    chunk_count = await process_document(db, document)
    return document, chunk_count


def _query_terms(query: str) -> list[str]:
    return [term.lower() for term in _WORD.findall(query) if len(term) > 1 and term.lower() not in _STOP_WORDS]


async def retrieve_document_sources(
    db: AsyncSession,
    *,
    query: str,
    document_ids: list[str],
    tenant_id: str,
    user_id: str,
    limit: int = 8,
    full_document: bool = False,
) -> list[WebSource]:
    if not document_ids:
        return []
    terms = _query_terms(query)
    # Targeted questions use PostgreSQL's GIN-backed full-text index. Keep the
    # Python scorer as the SQLite/test fallback and for whole-document tasks,
    # where ordered completeness matters more than rank.
    try:
        dialect_name = db.get_bind().dialect.name
    except (AttributeError, NotImplementedError):
        dialect_name = ""
    if not full_document and terms and dialect_name == "postgresql":
        # websearch_to_tsquery treats whitespace as AND. Retrieval should keep
        # useful chunks that match any meaningful query term, then rank chunks
        # matching more terms higher, mirroring the previous Python scorer.
        search_query = " OR ".join(terms)
        result = await db.execute(text("""
            SELECT c.text, c.ordinal, c.location, d.filename, d.id AS document_id,
                   ts_rank_cd(c.search_vector, websearch_to_tsquery('english', :query)) AS score
            FROM workspace_document_chunks c
            JOIN workspace_documents d ON d.id = c.document_id
            WHERE d.id = ANY(:document_ids)
              AND d.tenant_id = :tenant_id
              AND d.user_id = :user_id
              AND d.status = 'READY'
              AND d.expires_at > CURRENT_TIMESTAMP
              AND c.search_vector @@ websearch_to_tsquery('english', :query)
            ORDER BY score DESC, c.ordinal ASC
            LIMIT :limit
        """), {
            "query": search_query, "document_ids": document_ids,
            "tenant_id": tenant_id, "user_id": user_id, "limit": limit,
        })
        return [
            WebSource(
                title=f"{row.filename} — {row.location}", url="", snippet=row.text,
                provider="uploaded_document", freshness="uploaded",
                source_id=row.document_id,
            )
            for row in result
        ]
    result = await db.execute(
        select(DocumentChunk, Document)
        .join(Document, Document.id == DocumentChunk.document_id)
        .where(
            Document.id.in_(document_ids),
            Document.tenant_id == tenant_id,
            Document.user_id == user_id,
            Document.status == "READY",
            Document.expires_at > datetime.now(timezone.utc),
        )
    )
    phrase = " ".join(terms)
    ranked: list[tuple[float, DocumentChunk, Document]] = []
    for chunk, document in result.all():
        lowered = chunk.text.lower()
        matched = sum(1 for term in set(terms) if term in lowered)
        frequency = sum(min(lowered.count(term), 4) for term in set(terms))
        score = matched * 3 + frequency
        if phrase and phrase in lowered:
            score += 8
        if full_document or not terms:
            score = 1
        if score > 0:
            ranked.append((score, chunk, document))
    ranked.sort(key=lambda item: (-item[0], item[1].ordinal))
    # Whole-document tasks must see every extracted block. Applying the
    # question-answering top-k here silently dropped all chunks after the
    # eighth one, which made a "full document" report incomplete. Prompt
    # bounding happens later, after deterministic analysis has consumed the
    # complete ordered source set.
    selected = ranked if full_document else ranked[:limit]
    return [
        WebSource(
            title=f"{document.filename} — {chunk.location}",
            url="",
            snippet=chunk.text,
            provider="uploaded_document",
            freshness="uploaded",
            source_id=document.id,
        )
        for _, chunk, document in selected
    ]


async def resolve_conversation_document_ids(
    db: AsyncSession, *, conversation_id: str | None, requested_ids: list[str],
    tenant_id: str, user_id: str,
) -> list[str]:
    """Persist explicit attachments and restore them on later chat turns."""
    if not conversation_id:
        return requested_ids
    if requested_ids:
        valid = await db.execute(select(Document.id).where(
            Document.id.in_(requested_ids), Document.tenant_id == tenant_id,
            Document.user_id == user_id, Document.status == "READY",
            Document.expires_at > datetime.now(timezone.utc),
        ))
        valid_ids = list(valid.scalars().all())
        existing = await db.execute(select(ConversationDocument.document_id).where(
            ConversationDocument.conversation_id == conversation_id,
            ConversationDocument.tenant_id == tenant_id,
            ConversationDocument.user_id == user_id,
        ))
        existing_ids = set(existing.scalars().all())
        db.add_all([
            ConversationDocument(
                tenant_id=tenant_id, user_id=user_id, conversation_id=conversation_id,
                document_id=document_id,
            )
            for document_id in valid_ids if document_id not in existing_ids
        ])
        if valid_ids:
            await db.commit()
        return valid_ids
    result = await db.execute(select(ConversationDocument.document_id).where(
        ConversationDocument.conversation_id == conversation_id,
        ConversationDocument.tenant_id == tenant_id,
        ConversationDocument.user_id == user_id,
    ).order_by(ConversationDocument.created_at))
    return list(result.scalars().all())


async def list_documents(db: AsyncSession, *, tenant_id: str, user_id: str) -> list[tuple[Document, int]]:
    result = await db.execute(
        select(Document, func.count(DocumentChunk.id))
        .outerjoin(DocumentChunk, DocumentChunk.document_id == Document.id)
        .where(Document.tenant_id == tenant_id, Document.user_id == user_id)
        .group_by(Document.id)
        .order_by(Document.created_at.desc())
    )
    return [(document, int(chunk_count)) for document, chunk_count in result.all()]


async def get_document(
    db: AsyncSession, *, document_id: str, tenant_id: str, user_id: str
) -> tuple[Document, int] | None:
    result = await db.execute(
        select(Document, func.count(DocumentChunk.id))
        .outerjoin(DocumentChunk, DocumentChunk.document_id == Document.id)
        .where(
            Document.id == document_id,
            Document.tenant_id == tenant_id,
            Document.user_id == user_id,
        )
        .group_by(Document.id)
    )
    row = result.one_or_none()
    return (row[0], int(row[1])) if row else None


async def delete_document(
    db: AsyncSession, *, document_id: str, tenant_id: str, user_id: str
) -> bool:
    result = await db.execute(
        select(Document).where(
            Document.id == document_id,
            Document.tenant_id == tenant_id,
            Document.user_id == user_id,
        )
    )
    document = result.scalar_one_or_none()
    if document is None:
        return False
    storage_path = (_STORAGE_ROOT.parents[1] / document.storage_path).resolve()
    storage_root = _STORAGE_ROOT.resolve()
    await db.execute(delete(Document).where(Document.id == document.id))
    await db.commit()
    if document.storage_provider == "supabase" and document.storage_bucket and document.storage_object_path:
        await delete_object(document.storage_bucket, document.storage_object_path)
        return True
    if storage_path.is_relative_to(storage_root):
        try:
            storage_path.unlink(missing_ok=True)
        except OSError:
            # Database deletion is authoritative; an unavailable filesystem
            # cleanup can be handled by periodic orphan cleanup.
            pass
    return True
