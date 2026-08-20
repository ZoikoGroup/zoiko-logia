"""
Text extraction from uploaded documents.

One function per format, all returning the same shape: a list of Segments,
where the locator is the human-readable place the text came from ("page 3",
"sheet 'Summary'"). Chunking happens afterwards in chunker.py — keeping
extraction and chunking apart is what lets tabular formats keep their header
row on every chunk without the PDF path knowing anything about spreadsheets.

Nothing here raises on a merely difficult file. A page that yields no text is
skipped; a file that yields no text AT ALL raises ExtractionError so the
uploader is told the truth rather than being handed a document that is
technically "ready" and retrieves nothing.
"""
from __future__ import annotations

import csv
import io
from dataclasses import dataclass


class ExtractionError(Exception):
    """The file could not be turned into usable text. The message reaches the
    uploader, so it must say what to do about it."""


@dataclass
class Segment:
    locator: str
    text: str
    # True for spreadsheet/CSV segments, where the first line is a header row
    # that must be repeated on every chunk cut from this segment. Without it a
    # chunk reads "Q3 | 4,200 | 1,180" and neither the model nor the reader can
    # tell which column is revenue and which is tax.
    tabular: bool = False


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx", ".csv", ".txt", ".md"}

# A page or slide with fewer than this many characters is treated as decoration
# (a title slide, a scan artefact) rather than content.
_MIN_SEGMENT_CHARS = 20

# Guards against a spreadsheet with a million empty-but-styled rows producing
# a multi-hundred-megabyte string in memory.
_MAX_SHEET_ROWS = 20_000
_MAX_CELL_CHARS = 500


def _clean(text: str) -> str:
    """Collapse the whitespace soup that PDF and DOCX extraction produces,
    without joining what were genuinely separate paragraphs."""
    lines = [" ".join(line.split()) for line in (text or "").splitlines()]
    out: list[str] = []
    for line in lines:
        if line or (out and out[-1]):
            out.append(line)
    return "\n".join(out).strip()


def _extract_pdf(data: bytes) -> list[Segment]:
    from pypdf import PdfReader

    try:
        reader = PdfReader(io.BytesIO(data))
    except Exception as exc:
        raise ExtractionError(
            f"This PDF could not be opened ({type(exc).__name__}). It may be corrupt."
        ) from exc

    if getattr(reader, "is_encrypted", False):
        # An empty user password is the common "protected but readable" case
        # and is worth one attempt before giving up.
        try:
            reader.decrypt("")
        except Exception:
            raise ExtractionError(
                "This PDF is password-protected. Remove the password and upload it again."
            )

    segments: list[Segment] = []
    for index, page in enumerate(reader.pages, start=1):
        try:
            text = _clean(page.extract_text() or "")
        except Exception:
            continue                      # one unreadable page must not lose the rest
        if len(text) >= _MIN_SEGMENT_CHARS:
            segments.append(Segment(locator=f"page {index}", text=text))

    if not segments:
        raise ExtractionError(
            "No text could be read from this PDF. It is most likely a scan or a "
            "photo of a document, which needs OCR - try a text-based PDF, or "
            "export the original to .docx."
        )
    return segments


def _extract_docx(data: bytes) -> list[Segment]:
    import docx

    try:
        document = docx.Document(io.BytesIO(data))
    except Exception as exc:
        raise ExtractionError(
            f"This Word document could not be opened ({type(exc).__name__})."
        ) from exc

    parts = [p.text for p in document.paragraphs if p.text and p.text.strip()]

    # Tables in a Word file are usually the numbers the question is about, so
    # they are pulled out as pipe-delimited rows rather than dropped.
    for table_index, table in enumerate(document.tables, start=1):
        rows = []
        for row in table.rows:
            cells = [" ".join(c.text.split())[:_MAX_CELL_CHARS] for c in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"[Table {table_index}]\n" + "\n".join(rows))

    text = _clean("\n\n".join(parts))
    if len(text) < _MIN_SEGMENT_CHARS:
        raise ExtractionError("This Word document appears to be empty.")
    return [Segment(locator="document", text=text)]


def _extract_xlsx(data: bytes) -> list[Segment]:
    import openpyxl

    try:
        # read_only is not optional: normal mode builds a full cell object
        # graph, which on a large workbook is hundreds of megabytes.
        # data_only returns the cached result of a formula rather than its
        # text, which is what the reader means by the value.
        workbook = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception as exc:
        raise ExtractionError(
            f"This spreadsheet could not be opened ({type(exc).__name__})."
        ) from exc

    segments: list[Segment] = []
    try:
        for sheet in workbook.worksheets:
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                if len(rows) >= _MAX_SHEET_ROWS:
                    rows.append(f"[truncated at {_MAX_SHEET_ROWS} rows]")
                    break
                cells = [
                    ("" if v is None else str(v).replace("\n", " ").strip()[:_MAX_CELL_CHARS])
                    for v in row
                ]
                while cells and not cells[-1]:
                    cells.pop()
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                segments.append(
                    Segment(locator=f"sheet {sheet.title!r}", text="\n".join(rows), tabular=True)
                )
    finally:
        workbook.close()

    if not segments:
        raise ExtractionError("This spreadsheet contains no readable cells.")
    return segments


def _extract_pptx(data: bytes) -> list[Segment]:
    from pptx import Presentation

    try:
        presentation = Presentation(io.BytesIO(data))
    except Exception as exc:
        raise ExtractionError(
            f"This presentation could not be opened ({type(exc).__name__})."
        ) from exc

    segments: list[Segment] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [" ".join(c.text.split())[:_MAX_CELL_CHARS] for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
        text = _clean("\n".join(parts))
        if len(text) >= _MIN_SEGMENT_CHARS:
            segments.append(Segment(locator=f"slide {index}", text=text))

    if not segments:
        raise ExtractionError(
            "No text could be read from this presentation - the slides may be images."
        )
    return segments


def _decode(data: bytes) -> str:
    """Decode text bytes without ever raising. Exports from Excel are routinely
    cp1252 or UTF-16, so a hard utf-8 decode would reject good files."""
    for encoding in ("utf-8-sig", "utf-8", "utf-16", "cp1252"):
        try:
            return data.decode(encoding)
        except (UnicodeDecodeError, UnicodeError):
            continue
    return data.decode("utf-8", errors="replace")


def _extract_csv(data: bytes) -> list[Segment]:
    text = _decode(data)
    try:
        dialect: type[csv.Dialect] | csv.Dialect = csv.Sniffer().sniff(
            text[:8192], delimiters=",;\t|"
        )
    except csv.Error:
        dialect = csv.excel               # a single-column file sniffs as nothing

    rows: list[str] = []
    for row in csv.reader(io.StringIO(text), dialect):
        if len(rows) >= _MAX_SHEET_ROWS:
            rows.append(f"[truncated at {_MAX_SHEET_ROWS} rows]")
            break
        cells = [c.replace("\n", " ").strip()[:_MAX_CELL_CHARS] for c in row]
        if any(cells):
            rows.append(" | ".join(cells))

    if not rows:
        raise ExtractionError("This CSV file is empty.")
    return [Segment(locator="rows", text="\n".join(rows), tabular=True)]


def _extract_text(data: bytes) -> list[Segment]:
    text = _clean(_decode(data))
    if len(text) < _MIN_SEGMENT_CHARS:
        raise ExtractionError("This file is empty or too short to be useful.")
    return [Segment(locator="document", text=text)]


_EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".pptx": _extract_pptx,
    ".csv": _extract_csv,
    ".txt": _extract_text,
    ".md": _extract_text,
}


def extract(data: bytes, extension: str) -> list[Segment]:
    """Text segments from an uploaded file. Raises ExtractionError with a
    message intended for the uploader when the file yields nothing usable."""
    extractor = _EXTRACTORS.get(extension.lower())
    if extractor is None:
        raise ExtractionError(f"{extension!r} files cannot be read yet.")
    return extractor(data)
