import abc
import os
import asyncio
from typing import Optional

class BaseParser(abc.ABC):
    @abc.abstractmethod
    async def parse_file(self, file_path: str) -> str:
        """Parses a file at the given path and returns structured markdown content."""
        pass

class LlamaParseAdapter(BaseParser):
    """Primary parser using the LlamaParse cloud API."""
    def __init__(self, api_key: Optional[str] = None):
        self.api_key = api_key or os.environ.get("LLAMA_CLOUD_API_KEY")

    async def parse_file(self, file_path: str) -> str:
        if not self.api_key:
            raise ValueError("LlamaParse API Key is missing. Please set the LLAMA_CLOUD_API_KEY environment variable.")
        
        # import here to avoid startup issues if not installed
        from llama_parse import LlamaParse
        
        parser = LlamaParse(
            api_key=self.api_key,
            result_type="markdown",
            verbose=False
        )
        
        # Run loading asynchronously
        loop = asyncio.get_event_loop()
        documents = await loop.run_in_executor(None, parser.load_data, file_path)
        return "\n\n".join([doc.text for doc in documents])

class DoclingParserAdapter(BaseParser):
    """Fallback local parser using Docling library (falls back to basic text extractor if not installed)."""
    async def parse_file(self, file_path: str) -> str:
        try:
            from docling.document_converter import DocumentConverter
            converter = DocumentConverter()
            
            loop = asyncio.get_event_loop()
            result = await loop.run_in_executor(None, converter.convert, file_path)
            return result.document.export_to_markdown()
        except Exception as e:
            # Fallback to a simpler extraction method if docling import/conversion fails
            return self._fallback_simple_extract(file_path)

    def _fallback_simple_extract(self, file_path: str) -> str:
        ext = os.path.splitext(file_path)[1].lower()
        filename = os.path.basename(file_path)
        
        try:
            if ext == ".pdf":
                import pypdf
                reader = pypdf.PdfReader(file_path)
                return "\n\n".join([f"## Page {i+1}\n\n{page.extract_text()}" for i, page in enumerate(reader.pages)])
            
            elif ext == ".docx":
                import docx
                doc = docx.Document(file_path)
                return "\n\n".join([para.text for para in doc.paragraphs if para.text.strip()])
                
            elif ext == ".pptx":
                import pptx
                prs = pptx.Presentation(file_path)
                text = []
                for i, slide in enumerate(prs.slides):
                    text.append(f"## Slide {i+1}")
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                return "\n\n".join(text)
                
            elif ext == ".xlsx":
                import openpyxl
                wb = openpyxl.load_workbook(file_path, data_only=True)
                text = []
                for sheet in wb.worksheets:
                    text.append(f"## Sheet: {sheet.title}")
                    for row in sheet.iter_rows(values_only=True):
                        row_text = " | ".join([str(cell) for cell in row if cell is not None])
                        if row_text:
                            text.append(row_text)
                return "\n\n".join(text)
        except Exception as e:
            return f"# Document: {filename}\n\n[Parsing error occurred: {str(e)}]"
            
        # Extremely basic fallback when no parsing libraries are present
        return f"# Document: {filename}\n\n[Fallback simple text extraction placeholder for {ext} file type.]"

def get_parser(use_fallback: bool = False) -> BaseParser:
    """Factory helper to obtain the preferred parser service."""
    if not use_fallback and os.environ.get("LLAMA_CLOUD_API_KEY"):
        return LlamaParseAdapter()
    return DoclingParserAdapter()
