import io

from docx import Document as DocxDocument
from openpyxl import Workbook
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from app.domains.kriton_workspace.documents import extract_document
from app.orchestration.websearch import WebSource, build_web_grounded_prompt
from app.orchestration.document_pipeline import (
    analyse_spreadsheet_sources,
    build_document_generation_prompt,
    plan_document_task,
)
from app.orchestration.service import _terminal_response_state
from app.orchestration.routing_matrix import ROUTE_LLM, ROUTE_REFUSAL
from app.domains.kriton_workspace.artifacts import _render
from app.domains.kriton_workspace.document_spec import build_document_spec
from app.domains.kriton_workspace import documents as document_service
from app.domains.kriton_workspace.models import ConversationDocument, Document, DocumentChunk
from app.domains.kriton_workspace.documents import ingest_document, resolve_conversation_document_ids, retrieve_document_sources
from sqlalchemy.ext.asyncio import async_sessionmaker, create_async_engine
import pytest


def test_docx_extraction_preserves_text_and_tables():
    document = DocxDocument()
    document.add_heading("Revenue recognition", level=1)
    document.add_paragraph("Revenue is recognised when control transfers to the customer.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Year"
    table.cell(0, 1).text = "Revenue"
    table.cell(1, 0).text = "2026"
    table.cell(1, 1).text = "125000"
    buffer = io.BytesIO()
    document.save(buffer)

    chunks = extract_document("policy.docx", buffer.getvalue())

    assert any("control transfers" in chunk.text for chunk in chunks)
    assert any("125000" in chunk.text and chunk.location == "Table 1" for chunk in chunks)


def test_xlsx_extraction_preserves_sheet_location_formula_and_value():
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "P&L"
    sheet.append(["Account", "Amount"])
    sheet.append(["Revenue", 125000])
    sheet.append(["Costs", 75000])
    sheet.append(["Profit", "=B2-B3"])
    buffer = io.BytesIO()
    workbook.save(buffer)

    chunks = extract_document("accounts.xlsx", buffer.getvalue())

    assert len(chunks) == 1
    assert chunks[0].location.startswith("'P&L'!")
    assert "Revenue\t125000" in chunks[0].text
    assert "=B2-B3" in chunks[0].text


def test_document_evidence_prompt_is_grounded_and_has_no_public_url():
    source = WebSource(
        title="Accounts.xlsx — 'P&L'!1:4",
        url="",
        snippet="Revenue\t125000",
        provider="uploaded_document",
        source_id="doc-123",
    )

    prompt = build_web_grounded_prompt("What is revenue?", [source])

    assert "using ONLY the numbered evidence sources" in prompt
    assert "Document ID: doc-123" in prompt
    assert "Revenue\t125000" in prompt


def test_empty_evidence_prompt_fails_closed():
    prompt = build_web_grounded_prompt("What is revenue?", [])

    assert "Do not answer from model knowledge" in prompt
    assert "reliable evidence could not be retrieved" in prompt


def test_empty_evidence_prompt_can_allow_general_knowledge_without_fake_sources():
    prompt = build_web_grounded_prompt(
        "What is GBP?", [], allow_general_knowledge=True,
    )

    assert "answer directly using established general knowledge" in prompt
    assert "Do not invent citations" in prompt
    assert "source scope" in prompt
    assert "ask the user to attach" not in prompt


def test_management_report_uses_full_document_retrieval():
    plan = plan_document_task(
        "Create a management report based only on this uploaded workbook.",
        has_documents=True,
    )

    assert plan.task_type == "document_generation"
    assert plan.retrieval_mode == "full_document"
    assert plan.response_mode == "chat_with_artifact"
    assert plan.output_format == "docx"


def test_plain_summary_uses_full_document_without_creating_artifact():
    plan = plan_document_task("Summarize it", has_documents=True)

    assert plan.task_type == "document_summary"
    assert plan.retrieval_mode == "full_document"
    assert plan.response_mode == "chat"
    assert plan.output_format == "chat"


def test_terminal_response_state_is_available_before_artifact_generation():
    assert _terminal_response_state(False) == ("answered", ROUTE_LLM)
    assert _terminal_response_state(True) == ("refused", ROUTE_REFUSAL)


def test_explicit_summary_document_still_creates_download():
    plan = plan_document_task(
        "Create a PDF summary document from the attached workbook.",
        has_documents=True,
    )

    assert plan.task_type == "document_generation"
    assert plan.response_mode == "chat_with_artifact"
    assert plan.output_format == "pdf"


def test_working_paper_automatically_selects_xlsx():
    plan = plan_document_task(
        "Create an audit working paper from this workbook.", has_documents=True
    )
    assert plan.output_format == "xlsx"


def test_board_report_automatically_selects_pdf():
    plan = plan_document_task(
        "Prepare a board report from this workbook.", has_documents=True
    )
    assert plan.output_format == "pdf"


def test_profit_and_loss_generation_uses_full_document_retrieval():
    plan = plan_document_task(
        "Generate a profit-and-loss summary using the financial data in this workbook.",
        has_documents=True,
    )

    assert plan.retrieval_mode == "full_document"


def test_spreadsheet_analysis_calculates_verified_kpis_without_summary_double_counting():
    sources = [
        WebSource(
            title="sales.xlsx — 'Sales_Data'!1:3",
            url="",
            source_id="doc-1",
            provider="uploaded_document",
            snippet=(
                "Date\tCustomer\tProduct\tRevenue\tCost\tGross Profit\n"
                "2026-01-01\tApex Corp\tAnnual License\t125000\t42000\t83000\n"
                "2026-01-02\tBeta Ltd\tSupport\t75000\t30000\t45000"
            ),
        ),
        WebSource(
            title="sales.xlsx — 'Summary'!1:2",
            url="",
            source_id="doc-1",
            provider="uploaded_document",
            snippet="Metric\tRevenue\nTotal\t200000",
        ),
    ]

    analysis = analyse_spreadsheet_sources(sources)

    assert analysis["rows_analysed"] == 2
    assert analysis["metrics"]["total_revenue"] == 200000
    assert analysis["metrics"]["total_cost"] == 72000
    assert analysis["metrics"]["gross_profit"] == 128000
    assert analysis["customer_revenue"]["Apex Corp"] == 125000


def test_document_generation_prompt_constrains_gemini_to_verified_data():
    source = WebSource(
        title="sales.xlsx — 'Sales_Data'!1:2",
        url="",
        source_id="doc-1",
        provider="uploaded_document",
        snippet="Customer\tRevenue\nApex Corp\t125000",
    )
    prompt = build_document_generation_prompt(
        "Create a management report.", [source], {"metrics": {"total_revenue": 125000}}
    )

    assert "Use ONLY the evidence and verified calculations" in prompt
    assert '"total_revenue": 125000' in prompt
    assert "Apex Corp\t125000" in prompt


def test_document_generation_prompt_has_bounded_context():
    sources = [
        WebSource(title=f"sheet-{index}", url="", snippet="x" * 100_000, source_id="doc-1")
        for index in range(3)
    ]
    prompt = build_document_generation_prompt("Create a report", sources, {})

    assert len(prompt) < 65_000
    assert all(f"sheet-{index}" in prompt for index in range(3))


def test_large_document_prompt_remains_bounded_and_reports_complete_coverage():
    sources = [
        WebSource(
            title=f"sales.xlsx — 'Sheet-{index}'!1:40",
            url="",
            snippet="Customer\tRevenue\nApex\t100\n" + ("x" * 10_000),
            source_id="doc-1",
            provider="uploaded_document",
        )
        for index in range(25)
    ]

    analysis = analyse_spreadsheet_sources(sources)
    prompt = build_document_generation_prompt("Create a management report", sources, analysis)

    assert analysis["coverage"]["chunks_processed"] == 25
    assert analysis["coverage"]["complete"] is True
    assert len(prompt) < 70_000
    assert "Sheet-24" in prompt


def test_output_format_aliases_select_modern_download_formats():
    cases = {
        "Create a PDF board report": "pdf",
        "Create a PowerPoint presentation": "pptx",
        "Create a PPT slide deck": "pptx",
        "Create an XLS management report": "xlsx",
        "Create an Excel report": "xlsx",
        "Create a Word document": "docx",
        "Create a DOC report": "docx",
        "Export a Markdown report": "md",
    }
    for query, expected in cases.items():
        assert plan_document_task(query, has_documents=True).output_format == expected


def test_every_supported_artifact_format_opens_with_its_native_parser():
    narrative = (
        "## Executive Summary\nRevenue increased.\n"
        "## KPIs\n- Total Revenue: 125000\n- Gross Profit: 83000\n"
        "## Customer Table\n| Customer | Revenue |\n| --- | ---: |\n| Apex | 125000 |\n"
        "## Recommendations\n- Monitor customer concentration."
    )
    analysis = {
        "metrics": {"total_revenue": 125000, "gross_profit": 83000},
        "customer_revenue": {"Apex": 125000},
        "product_revenue": {"Software": 125000},
        "evidence_locations": ["sales.xlsx — Sales!1:2"],
    }

    request = "Create a report with a chart and flow diagram"
    spec = build_document_spec("Management Report", narrative, analysis, request)
    assert {block.type for block in spec.blocks} >= {"table", "chart", "flow_diagram"}

    docx_bytes, extension, _ = _render("docx", "Management Report", narrative, analysis, request)
    assert extension == "docx"
    word = DocxDocument(io.BytesIO(docx_bytes))
    assert word.paragraphs
    assert word.tables
    assert len(word.inline_shapes) >= 2

    xlsx_bytes, extension, _ = _render("xlsx", "Management Report", narrative, analysis, request)
    assert extension == "xlsx"
    workbook = load_workbook(io.BytesIO(xlsx_bytes), read_only=True)
    try:
        assert "Verified KPIs" in workbook.sheetnames
    finally:
        workbook.close()

    pdf_bytes, extension, _ = _render("pdf", "Management Report", narrative, analysis, request)
    assert extension == "pdf"
    assert len(PdfReader(io.BytesIO(pdf_bytes)).pages) >= 1

    pptx_bytes, extension, _ = _render("pptx", "Management Report", narrative, analysis, request)
    assert extension == "pptx"
    slides = Presentation(io.BytesIO(pptx_bytes)).slides
    assert len(slides) >= 4
    assert any(shape.has_table for slide in slides for shape in slide.shapes)
    assert sum(shape.shape_type == 13 for slide in slides for shape in slide.shapes) >= 2

    md_bytes, extension, mime_type = _render("md", "Management Report", narrative, analysis, request)
    assert extension == "md"
    assert mime_type.startswith("text/markdown")
    assert md_bytes.decode("utf-8").startswith("# Management Report")
    assert "| Customer | Revenue |" in md_bytes.decode("utf-8")
    assert md_bytes.decode("utf-8").count("data:image/png;base64,") >= 2


def test_document_spec_rejects_unresolved_financial_placeholders():
    with pytest.raises(ValueError, match="unresolved placeholder"):
        build_document_spec(
            "Management Report",
            "## Recommendation\nPending deal value is USD ?.",
            {},
        )


def test_document_exports_normalize_unicode_and_remove_internal_markers():
    narrative = (
        "## Executive Summary\n"
        "High\u2011value customers delivered 51\u202f% of revenue [REF\u20112].\n"
        "---\n"
        "## KPI Overview\n"
        "| KPI | Value |\n| --- | ---: |\n| In\u2011Progress Deals | 4 |\n"
        "Source: Summary sheet ([REF-2])."
    )
    spec = build_document_spec(
        "Kriton\u00a0Management Report", narrative,
        {"evidence_locations": ["sales.xlsx \u2014 Summary!1:2"]},
    )

    searchable = " ".join(
        [spec.title]
        + [block.text for block in spec.blocks]
        + [cell for block in spec.blocks for row in block.rows for cell in row]
    )
    assert "Kriton Management Report" == spec.title
    assert "High-value" in searchable
    assert "In-Progress" in searchable
    assert "REF" not in searchable
    assert not any(block.type == "paragraph" and block.text == "---" for block in spec.blocks)

    pdf_bytes, _, _ = _render("pdf", spec.title, narrative, {"evidence_locations": []})
    pdf_text = "\n".join(page.extract_text() or "" for page in PdfReader(io.BytesIO(pdf_bytes)).pages)
    assert "High-value" in pdf_text
    assert "In-Progress" in pdf_text
    assert "REF" not in pdf_text


@pytest.mark.asyncio
async def test_upload_to_conversation_retrieval_end_to_end(monkeypatch, tmp_path):
    """Exercise the real XLSX upload, extraction, DB indexing and retrieval path."""
    storage_root = tmp_path / "backend" / "data" / "workspace_documents"
    monkeypatch.setattr(document_service, "_STORAGE_ROOT", storage_root)
    monkeypatch.setattr(document_service, "uses_supabase", lambda: False)

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Sales_Data"
    sheet.append(["Customer", "Revenue"])
    sheet.append(["Summit Capital", 205000])
    sheet.append(["Apex Corp", 125000])
    content = io.BytesIO()
    workbook.save(content)

    engine = create_async_engine("sqlite+aiosqlite:///:memory:")
    async with engine.begin() as connection:
        for table in (Document.__table__, DocumentChunk.__table__, ConversationDocument.__table__):
            await connection.run_sync(table.create)
    sessions = async_sessionmaker(engine, expire_on_commit=False)
    try:
        async with sessions() as db:
            document, chunk_count = await ingest_document(
                db,
                filename="sales.xlsx",
                mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                content=content.getvalue(),
                tenant_id="tenant-1",
                user_id="user-1",
            )
            assert document.status == "READY"
            assert chunk_count == 1

            attached_ids = await resolve_conversation_document_ids(
                db,
                conversation_id="conversation-1",
                requested_ids=[document.id],
                tenant_id="tenant-1",
                user_id="user-1",
            )
            sources = await retrieve_document_sources(
                db,
                query="Which customer generated the most revenue?",
                document_ids=attached_ids,
                tenant_id="tenant-1",
                user_id="user-1",
            )
            restored_ids = await resolve_conversation_document_ids(
                db,
                conversation_id="conversation-1",
                requested_ids=[],
                tenant_id="tenant-1",
                user_id="user-1",
            )

            assert restored_ids == [document.id]
            assert len(sources) == 1
            assert "Summit Capital\t205000" in sources[0].snippet
            assert sources[0].source_id == document.id
    finally:
        await engine.dispose()
